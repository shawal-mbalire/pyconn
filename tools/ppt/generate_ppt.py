"""
tools/ppt/generate_ppt.py — Generate project architecture slide deck.

Produces a .pptx covering:
  - System overview & hardware pairing
  - Four core patterns (Static Heap, IRQ Split, Watchdog, Parser Pipeline)
  - Monorepo layout
  - Firmware data-flow diagram (text art)

Usage
-----
    pip install python-pptx
    python tools/ppt/generate_ppt.py            # -> output/pyconn_overview.pptx
    python tools/ppt/generate_ppt.py --out my.pptx
"""
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
_BG      = RGBColor(0x0D, 0x11, 0x17)   # near-black
_ACCENT  = RGBColor(0x00, 0xD4, 0x8A)   # teal-green
_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
_GREY    = RGBColor(0x8B, 0x94, 0x9E)
_YELLOW  = RGBColor(0xFF, 0xD7, 0x00)

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
_W = Inches(13.33)
_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def _bg(slide):
    """Fill slide background with dark colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _BG


def _add_textbox(slide, text, left, top, width, height,
                 size=20, bold=False, color=_WHITE, align=PP_ALIGN.LEFT,
                 wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def _accent_bar(slide, top=Inches(1.05)):
    """Thin horizontal accent rule under the title."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), top, Inches(12.33), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _ACCENT
    line.line.fill.background()


def _title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide)
    _add_textbox(slide, "pyconn",
                 Inches(0.5), Inches(1.8), Inches(12), Inches(1.2),
                 size=54, bold=True, color=_ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide,
                 "Hardened MicroPython  |  ESP32-C3  +  60GHz mmWave Radar",
                 Inches(0.5), Inches(3.1), Inches(12), Inches(0.6),
                 size=22, color=_WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide,
                 "Python-native RTOS  •  Zero-alloc hot path  •  Cooperative WDT",
                 Inches(0.5), Inches(3.8), Inches(12), Inches(0.5),
                 size=16, color=_GREY, align=PP_ALIGN.CENTER)
    _accent_bar(slide, top=Inches(4.6))
    _add_textbox(slide, "github.com/shawal-mbalire/pyconn",
                 Inches(0.5), Inches(4.8), Inches(12), Inches(0.4),
                 size=13, color=_GREY, align=PP_ALIGN.CENTER)


def _pattern_slide(prs, number, title, subtitle, bullets, code_snippet=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _add_textbox(slide, f"Pattern {number}", Inches(0.5), Inches(0.2),
                 Inches(4), Inches(0.4), size=12, color=_ACCENT)
    _add_textbox(slide, title, Inches(0.5), Inches(0.55),
                 Inches(12), Inches(0.65), size=30, bold=True, color=_WHITE)
    _add_textbox(slide, subtitle, Inches(0.5), Inches(1.25),
                 Inches(12), Inches(0.4), size=14, color=_GREY)
    _accent_bar(slide, top=Inches(1.7))

    # Bullet points
    bullet_text = "\n".join(f"  {b}" for b in bullets)
    _add_textbox(slide, bullet_text, Inches(0.5), Inches(1.85),
                 Inches(5.8), Inches(4.5), size=14, color=_WHITE)

    # Code block (right panel)
    if code_snippet:
        box = slide.shapes.add_shape(1,
            Inches(6.6), Inches(1.85), Inches(6.2), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
        box.line.color.rgb = _ACCENT

        txBox = slide.shapes.add_textbox(
            Inches(6.75), Inches(2.0), Inches(5.9), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = False
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = code_snippet
        run.font.size = Pt(10)
        run.font.color.rgb = _ACCENT
        run.font.name = "Courier New"


def _layout_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _add_textbox(slide, "Monorepo Layout", Inches(0.5), Inches(0.2),
                 Inches(12), Inches(0.65), size=30, bold=True, color=_WHITE)
    _accent_bar(slide, top=Inches(0.95))

    tree = (
        "/firmware\n"
        "  main.py          # Entry point + WDT init\n"
        "  boot.py          # Static heap pre-allocation\n"
        "  core/\n"
        "    scheduler.py   # uasyncio + cooperative heartbeats\n"
        "    memory.py      # GC management + BufferPool\n"
        "  drivers/\n"
        "    mr60_radar.py  # IRQ -> schedule -> state machine\n"
        "  app/\n"
        "    processor.py   # Frame consumer -> event dict\n"
        "/docs              # mmWave protocol + state machine specs\n"
        "/tools\n"
        "  obfuscate.py     # .py -> .mpy + name mangling\n"
        "  ppt/             # This slide generator\n"
        "justfile           # 11 top-level recipes"
    )
    box = slide.shapes.add_shape(1,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
    box.line.color.rgb = _ACCENT

    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.25), Inches(12.0), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    run = p.add_run()
    run.text = tree
    run.font.size = Pt(13)
    run.font.color.rgb = _ACCENT
    run.font.name = "Courier New"


def _dataflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _add_textbox(slide, "Firmware Data-Flow",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.65),
                 size=30, bold=True, color=_WHITE)
    _accent_bar(slide, top=Inches(0.95))

    flow = (
        "  MR60 Radar (115200 baud UART)\n"
        "       |\n"
        "  [UART HW rxbuf 2 KB]   <-- catches bytes during GC pauses\n"
        "       |\n"
        "  uart.irq(hard=False)   <-- fires between VM bytecodes\n"
        "       |\n"
        "  micropython.schedule(_drain)\n"
        "       |\n"
        "  _drain() -> readinto(RADAR_RX_BUF)   <-- zero new allocs\n"
        "       |\n"
        "  State Machine (_feed byte-by-byte)\n"
        "    WAIT_H1 -> WAIT_H2 -> READ_LEN -> READ_PAY -> VERIFY\n"
        "       |\n"
        "  @micropython.native _checksum()   <-- RISC-V native code\n"
        "       |\n"
        "  scheduler.enqueue(memoryview, length)   <-- zero-copy\n"
        "       |\n"
        "  Processor.run()  -- async consumer\n"
        "       |\n"
        "  on_event(dict)  -->  MQTT / ESP-NOW / AI inference"
    )

    box = slide.shapes.add_shape(1,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
    box.line.color.rgb = _ACCENT

    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.25), Inches(12.0), Inches(5.6))
    tf = txBox.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    run = p.add_run()
    run.text = flow
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x79, 0xC0, 0xFF)
    run.font.name = "Courier New"


# ── Slide content ─────────────────────────────────────────────────────────────

_PATTERNS = [
    {
        "title":    "Pre-allocated Static Heap",
        "subtitle": "Prevent GC fragmentation with fixed-size bytearrays at boot",
        "bullets": [
            "Allocate RADAR_RX_BUF (2 KB) and RADAR_FRAME_BUF (256 B) in boot.py",
            "Write into existing buffers — never create new bytearray in the loop",
            "Use memoryview for zero-copy slicing (no bytes() in hot path)",
            "Keeps heap 'flat' -> GC pauses < 1 ms vs. 10-20 ms fragmented",
            "boot.maybe_collect() called only from scheduler tick (safe point)",
        ],
        "code": (
            "# boot.py\n"
            "micropython.alloc_emergency_exception_buf(100)\n\n"
            "RADAR_RX_BUF   = bytearray(2048)\n"
            "RADAR_FRAME_BUF = bytearray(256)\n\n"
            "gc.collect()\n"
            "HEAP_BASELINE_FREE = gc.mem_free()\n\n"
            "# driver — zero new allocs\n"
            "n = uart.readinto(boot.RADAR_RX_BUF)\n"
            "mv = memoryview(boot.RADAR_FRAME_BUF)"
        ),
    },
    {
        "title":    "Hard/Soft Interrupt Split",
        "subtitle": "Near-instant hardware response with full Python safety",
        "bullets": [
            "Hard IRQ: fires at hardware speed, memory allocation FORBIDDEN",
            "Soft IRQ (hard=False): runs between VM bytecodes — full Python OK",
            "IRQ does ONE thing: micropython.schedule(drain_fn, 0)",
            "_drain_pending flag prevents duplicate schedules",
            "Result: < 1 bytecode delay for response, no alloc restriction",
        ],
        "code": (
            "# drivers/mr60_radar.py\n"
            "def _uart_irq(self, uart):\n"
            "    if not self._drain_pending:\n"
            "        self._drain_pending = True\n"
            "        micropython.schedule(\n"
            "            self._drain, 0\n"
            "        )\n\n"
            "def _drain(self, _):\n"
            "    self._drain_pending = False\n"
            "    n = self._uart.readinto(\n"
            "        self._rx_buf\n"
            "    )\n"
            "    if n:\n"
            "        for i in range(n):\n"
            "            self._feed(self._rx_buf[i])"
        ),
    },
    {
        "title":    "Cooperative Heartbeat Watchdog",
        "subtitle": "WDT is fed only when every task proves it is alive",
        "bullets": [
            "Each task owns a named slot in Scheduler._heartbeats dict",
            "task.heartbeat('name') updates ticks_ms timestamp each iteration",
            "Scheduler checks ALL heartbeats every second before wdt.feed()",
            "If any task misses its 5 s deadline -> WDT starves -> reboot",
            "Prevents 'alive-but-broken': MQTT hang can no longer hide radar failure",
        ],
        "code": (
            "# core/scheduler.py\n"
            "HEARTBEAT_TIMEOUT_MS = 5_000\n\n"
            "def _all_healthy(self):\n"
            "    now = time.ticks_ms()\n"
            "    for name, last in \\\n"
            "            self._heartbeats.items():\n"
            "        diff = time.ticks_diff(\n"
            "            now, last)\n"
            "        if diff > \\\n"
            "           self.HEARTBEAT_TIMEOUT_MS:\n"
            "            return False, name\n"
            "    return True, None\n\n"
            "async def run(self):\n"
            "    while True:\n"
            "        ok, _ = self._all_healthy()\n"
            "        if ok:\n"
            "            self._wdt.feed()"
        ),
    },
    {
        "title":    "State-Machine Frame Parser",
        "subtitle": "Non-blocking byte-by-byte parsing of MR60 radar frames",
        "bullets": [
            "Five states: WAIT_H1 -> WAIT_H2 -> READ_LEN -> READ_PAY -> VERIFY",
            "Never calls uart.read(n) — no blocking, no missed packets",
            "@micropython.native checksum -> compiled to RISC-V machine code",
            "on_frame() passes memoryview slice — zero bytes copied",
            "Bad checksum: silent discard + resync to WAIT_H1 (no crash)",
        ],
        "code": (
            "# @micropython.native = RISC-V codegen\n"
            "@micropython.native\n"
            "def _checksum(data, length):\n"
            "    s = 0\n"
            "    for i in range(length):\n"
            "        s = (s + data[i]) & 0xFF\n"
            "    return s\n\n"
            "# State machine (hot path)\n"
            "elif s == _S_VERIFY:\n"
            "    if byte == _checksum(\n"
            "        self._frame_buf,\n"
            "        self._payload_len\n"
            "    ):\n"
            "        self._on_frame(\n"
            "            self._mv[:self._payload_len],\n"
            "            self._payload_len\n"
            "        )\n"
            "    self._state = _S_WAIT_H1"
        ),
    },
]


def build(out_path: str) -> None:
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H

    _title_slide(prs)

    for i, pat in enumerate(_PATTERNS, start=1):
        _pattern_slide(
            prs, i,
            pat["title"],
            pat["subtitle"],
            pat["bullets"],
            pat.get("code", ""),
        )

    _layout_slide(prs)
    _dataflow_slide(prs)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate pyconn slide deck")
    parser.add_argument("--out", default="output/pyconn_overview.pptx",
                        help="Output .pptx path")
    args = parser.parse_args()
    build(args.out)
